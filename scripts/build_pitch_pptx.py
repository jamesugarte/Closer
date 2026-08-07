#!/usr/bin/env python3
"""Generate Closer Infosys AI case PowerPoint deck."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# Palette — bank-interview clean (avoid purple/cream AI defaults)
BG = RGBColor(0x0F, 0x1C, 0x1A)
SURFACE = RGBColor(0x1A, 0x2B, 0x28)
ACCENT = RGBColor(0x2F, 0x9E, 0x7A)
ACCENT_SOFT = RGBColor(0x3D, 0xB8, 0x8E)
WHITE = RGBColor(0xF7, 0xF5, 0xF0)
MUTED = RGBColor(0xA8, 0xB5, 0xB0)
WARN = RGBColor(0xE8, 0xA8, 0x3C)
DANGER = RGBColor(0xE0, 0x6C, 0x5C)
CARD = RGBColor(0x14, 0x24, 0x21)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_run(run, size=18, bold=False, color=WHITE, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_para(tf, text, size=16, bold=False, color=WHITE, space_before=6):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return p


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bg(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    fill_shape(shape, BG)


def add_accent_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), SLIDE_H)
    fill_shape(bar, ACCENT)


def add_card(slide, left, top, width, height):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_shape(card, CARD)
    card.adjustments[0] = 0.08
    return card


def add_footer(slide, page, total=5):
    add_textbox(
        slide,
        Inches(0.5),
        Inches(7.05),
        Inches(10),
        Inches(0.35),
        "Infosys Consulting  ·  AI Case Presentation  ·  FS Banking & Payments",
        size=11,
        color=MUTED,
    )
    add_textbox(
        slide,
        Inches(11.5),
        Inches(7.05),
        Inches(1.5),
        Inches(0.35),
        f"{page} / {total}",
        size=11,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)

    add_textbox(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.35),
                "CANDIDATE BRIEF  ·  AI-LED PROCESS IN BANKING", size=12, bold=True, color=ACCENT_SOFT)

    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.9),
                "Closer", size=48, bold=True, color=WHITE)

    add_textbox(slide, Inches(0.6), Inches(1.85), Inches(11.5), Inches(0.9),
                "An AI banking experience for students who finally have money —\nand no idea what they can safely spend.",
                size=22, color=MUTED)

    # Two cards
    add_card(slide, Inches(0.6), Inches(3.1), Inches(5.8), Inches(2.5))
    add_textbox(slide, Inches(0.85), Inches(3.25), Inches(5.3), Inches(0.35),
                "AUDIENCE (FROM BRIEF)", size=12, bold=True, color=ACCENT_SOFT)
    box = slide.shapes.add_textbox(Inches(0.85), Inches(3.7), Inches(5.3), Inches(1.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "College students managing money independently for the first time."
    set_run(run, size=16, color=WHITE)
    add_para(tf, "Income from jobs / internships · tuition & rent · multiple payment apps (Venmo, campus card, checking).", size=14, color=MUTED, space_before=10)

    add_card(slide, Inches(6.8), Inches(3.1), Inches(5.8), Inches(2.5))
    add_textbox(slide, Inches(7.05), Inches(3.25), Inches(5.3), Inches(0.35),
                "CORE PROBLEM (FROM BRIEF)", size=12, bold=True, color=ACCENT_SOFT)
    box = slide.shapes.add_textbox(Inches(7.05), Inches(3.7), Inches(5.3), Inches(1.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "They struggle to understand how much they can safely spend."
    set_run(run, size=16, bold=True, color=WHITE)
    add_para(tf, "Not “teach accounting.” Make safe spend, shocks, and goals legible — then act with AI.", size=14, color=MUTED, space_before=10)

    add_textbox(slide, Inches(0.6), Inches(5.85), Inches(12), Inches(0.7),
                "POV: Safe spend = balance after protected bills, risk cushion, and goal earmarks. Progress = days closer to real goals.",
                size=15, color=WHITE)

    add_footer(slide, 1)


def slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)

    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.3),
                "KEY OBSERVATIONS", size=12, bold=True, color=ACCENT_SOFT)
    add_textbox(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.55),
                "Students are overconfident — and underinsured", size=28, bold=True, color=WHITE)

    # Stat cards
    stats = [
        ("64%", "Confident managing\nbasic budgeting & saving", ACCENT_SOFT, "CFP Board · Fall 2025"),
        ("~25%", "Solidly understand full\ncollege cost to budget", WARN, "IHE Student Voice · 2025"),
        ("42%", "Gen Z living paycheck\nto paycheck", DANGER, "Bank of America · 2026"),
    ]
    x = 0.6
    for val, label, color, src in stats:
        add_card(slide, Inches(x), Inches(1.45), Inches(3.9), Inches(1.85))
        add_textbox(slide, Inches(x + 0.25), Inches(1.6), Inches(3.4), Inches(0.6),
                    val, size=36, bold=True, color=color)
        add_textbox(slide, Inches(x + 0.25), Inches(2.25), Inches(3.4), Inches(0.7),
                    label, size=14, color=WHITE)
        add_textbox(slide, Inches(x + 0.25), Inches(2.95), Inches(3.4), Inches(0.25),
                    src, size=10, color=MUTED)
        x += 4.15

    # Fragility bars (simple visual)
    add_textbox(slide, Inches(0.6), Inches(3.5), Inches(12), Inches(0.35),
                "FRAGILITY — THEY CANNOT PLAN FOR THE UNEXPECTED", size=12, bold=True, color=ACCENT_SOFT)

    bars = [
        ("Trouble raising $500 for a shock", 56, "Trellis SFWS Fall 2024"),
        ("Ran out of money at least once (2024)", 68, "Trellis SFWS Fall 2024"),
        ("≤$1,000 shock could threaten enrollment", 36, "IHE Student Voice 2025"),
        ("Unaware if campus offers emergency aid", 64, "IHE Student Voice 2025"),
    ]
    y = 3.9
    for label, pct, _ in bars:
        add_textbox(slide, Inches(0.6), Inches(y), Inches(5.2), Inches(0.28),
                    label, size=13, color=WHITE)
        track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(y + 0.05), Inches(5.5), Inches(0.2))
        fill_shape(track, SURFACE)
        track.adjustments[0] = 0.5
        fill_w = 5.5 * (pct / 100)
        fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(y + 0.05), Inches(fill_w), Inches(0.2))
        fill_shape(fill, ACCENT if pct < 60 else WARN)
        fill.adjustments[0] = 0.5
        add_textbox(slide, Inches(11.5), Inches(y), Inches(1.2), Inches(0.28),
                    f"{pct}%", size=13, bold=True, color=WHITE)
        y += 0.42

    add_textbox(slide, Inches(0.6), Inches(5.75), Inches(12), Inches(0.9),
                "Also: Convenience waste (delivery, rideshare, subscriptions) is where “I’ll save later” dies — Gen Z cuts apparel but still spends on experiences (McKinsey / PwC). "
                "89% of student card users put basics on plastic (Trellis).",
                size=13, color=MUTED)

    add_footer(slide, 2)


def slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)

    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.3),
                "RECOMMENDED APPROACH", size=12, bold=True, color=ACCENT_SOFT)
    add_textbox(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.55),
                "Measure saving in units of time", size=28, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.4),
                "Dollars are abstract. “Six days earlier” is felt.", size=16, color=MUTED)

    add_card(slide, Inches(0.6), Inches(1.9), Inches(5.8), Inches(2.4))
    add_textbox(slide, Inches(0.85), Inches(2.05), Inches(5.3), Inches(0.3),
                "WHAT BANKS USUALLY SHIP", size=12, bold=True, color=MUTED)
    box = slide.shapes.add_textbox(Inches(0.85), Inches(2.45), Inches(5.3), Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "“You spent $14 on DoorDash.”"
    set_run(run, size=16, color=WHITE)
    add_para(tf, "Shame → ignore → repeat.", size=14, color=MUTED, space_before=6)
    add_para(tf, "Category pies with no calendar — no link to concert night or spring break.", size=14, color=MUTED, space_before=10)

    add_card(slide, Inches(6.8), Inches(1.9), Inches(5.8), Inches(2.4))
    add_textbox(slide, Inches(7.05), Inches(2.05), Inches(5.3), Inches(0.3),
                "WHAT WE RECOMMEND", size=12, bold=True, color=ACCENT_SOFT)
    box = slide.shapes.add_textbox(Inches(7.05), Inches(2.45), Inches(5.3), Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "“Skip one delivery → Fall concert moves 5 days closer.”"
    set_run(run, size=16, bold=True, color=WHITE)
    add_para(tf, "Same $14 — now it’s a date on the wall.", size=14, color=MUTED, space_before=6)
    add_para(tf, "AI tips ranked by days gained vs lifestyle friction. Student accepts or rejects.", size=14, color=MUTED, space_before=10)

    pillars = [
        ("1 · SAFE SPEND", "Free-to-spend = checking − protected bills (rent, tuition, meal plan, car loan) − risk cushion − goal earmarks."),
        ("2 · TIME AS CURRENCY", "Every AI action answers: how many days closer? Habit cuts, surplus moves, and portfolio reshuffles share that score."),
        ("3 · HONEST PORTFOLIOS", "Fixed dates (concert, spring break) don’t slide. Flexible wants can yield. Obligations stay obligations — never fake goals."),
    ]
    x = 0.6
    for title, body in pillars:
        add_card(slide, Inches(x), Inches(4.55), Inches(3.9), Inches(2.0))
        add_textbox(slide, Inches(x + 0.2), Inches(4.7), Inches(3.5), Inches(0.35),
                    title, size=12, bold=True, color=ACCENT_SOFT)
        add_textbox(slide, Inches(x + 0.2), Inches(5.15), Inches(3.5), Inches(1.2),
                    body, size=13, color=WHITE)
        x += 4.15

    add_footer(slide, 3)


def slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)

    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.3),
                "AI-BUILT PROTOTYPE", size=12, bold=True, color=ACCENT_SOFT)
    add_textbox(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.55),
                "Closer — live decisioning on a student ledger", size=28, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.35),
                "Vibe-coded Next.js demo · mock multi-rail history · two personas", size=15, color=MUTED)

    add_card(slide, Inches(0.6), Inches(1.85), Inches(5.8), Inches(2.0))
    add_textbox(slide, Inches(0.85), Inches(2.0), Inches(5.3), Inches(0.3),
                "JORDAN · COURSE CORRECT", size=13, bold=True, color=WARN)
    add_textbox(slide, Inches(0.85), Inches(2.4), Inches(5.3), Inches(1.2),
                "Freshman cash-flow fire drill — delivery, rideshare, BNPL, thin risk cushion. Diagnostic + habit tips that pull a fantasy goal back toward reality.",
                size=14, color=WHITE)

    add_card(slide, Inches(6.8), Inches(1.85), Inches(5.8), Inches(2.0))
    add_textbox(slide, Inches(7.05), Inches(2.0), Inches(5.3), Inches(0.3),
                "MAYA · PORTFOLIO SQUEEZE", size=13, bold=True, color=ACCENT_SOFT)
    add_textbox(slide, Inches(7.05), Inches(2.4), Inches(5.3), Inches(1.2),
                "“Doing fine” sophomore — car loan on protected bills (not a goal), competing wants (concert, AirPods, spring break). Freedom + security without austerity theater.",
                size=14, color=WHITE)

    add_textbox(slide, Inches(0.6), Inches(4.1), Inches(12), Inches(0.3),
                "WHAT THE AI ACTUALLY DOES", size=12, bold=True, color=ACCENT_SOFT)

    bullets = [
        "Surfaces unnecessary spend patterns (delivery, rideshare, BNPL, coffee) as days-closer tips",
        "Moves free-to-spend into goals; rearranges reserves across ranked goals",
        "Protects fixed deadlines; never slides spring break to fund a soft want",
        "Keeps car loan / rent / tuition as protected obligations — not savings goals",
        "Student agency: accept or reject every recommendation",
    ]
    y = 4.45
    for b in bullets:
        add_textbox(slide, Inches(0.6), Inches(y), Inches(12), Inches(0.32),
                    f"→  {b}", size=14, color=WHITE)
        y += 0.32

    add_textbox(slide, Inches(0.6), Inches(6.15), Inches(12), Inches(0.55),
                "Banking insight: AI sits on deposit + payment rails students already use. Value is decisioning — safe spend, shock buffers, goal pacing — not another wallet.",
                size=13, color=MUTED)

    add_footer(slide, 4)


def slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)

    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.3),
                "DEMO PATH  ·  INSIGHTS & CONSIDERATIONS", size=12, bold=True, color=ACCENT_SOFT)
    add_textbox(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.55),
                "Live walkthrough · Maya (~4 minutes)", size=28, bold=True, color=WHITE)

    steps = [
        ("01", "Login → diagnostic", "Show car loan inside protected obligations. Free-to-spend after bills + risk + earmarks."),
        ("02", "Goals · ranked feasibility", "Expand fixed-date goals vs flexible AirPods. Portfolio optimize tips below."),
        ("03", "Accept an AI tip", "Rearrange or surplus → days closer. Bills untouched. Reject shows agency."),
        ("04", "Add ~$2,000 goal as #1", "Tips refresh without sliding hard dates or raiding the car note. Optional: Jordan rescue arc."),
    ]
    y = 1.45
    for num, title, body in steps:
        add_card(slide, Inches(0.6), Inches(y), Inches(12.1), Inches(0.85))
        add_textbox(slide, Inches(0.85), Inches(y + 0.15), Inches(0.7), Inches(0.5),
                    num, size=20, bold=True, color=ACCENT_SOFT)
        add_textbox(slide, Inches(1.7), Inches(y + 0.12), Inches(10.5), Inches(0.3),
                    title, size=16, bold=True, color=WHITE)
        add_textbox(slide, Inches(1.7), Inches(y + 0.45), Inches(10.5), Inches(0.3),
                    body, size=13, color=MUTED)
        y += 0.95

    add_textbox(slide, Inches(0.6), Inches(5.4), Inches(12), Inches(0.9),
                "Close: Closer answers the brief — how much can I safely spend? — then turns leftover life into days toward goals, while keeping students enrolled through shocks they don’t plan for.",
                size=15, bold=True, color=WHITE)

    add_textbox(slide, Inches(0.6), Inches(6.3), Inches(12), Inches(0.45),
                "Considerations: time framing > dollar shaming · fixed-deadline product rules (not just LLM copy) · next: real aggregators, model risk, explainability for regulated advice.",
                size=12, color=MUTED)

    add_footer(slide, 5)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)

    out = "/Users/jamesugarte/Projects/closer/Closer_Infosys_AI_Case_Deck.pptx"
    prs.save(out)
    print(out)


if __name__ == "__main__":
    main()
